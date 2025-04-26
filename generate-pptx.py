from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
import os
import tkinter.filedialog

# Remove the existing presentation file if it exists
if os.path.exists("../本番スライド/2024字幕.pptx"):
    os.remove("../本番スライド/2024字幕.pptx")

# Path to the text files
es_path = './24es_text.txt'
ja_path = './24ja_text.txt'

# Read the content of the files
with open(es_path, 'r', encoding='utf-8') as file:
    es_text = file.readlines()

with open(ja_path, 'r', encoding='utf-8') as file:
    ja_text = file.readlines()



def specifySaveLocation():
    global output_path
    # Open a file dialog to select the output path
    output_path = tkinter.filedialog.askdirectory()
    if not output_path:
        print("No directory selected. Using default path.")
        output_path = '.'  # Default path if no directory is selected
    return output_path


# Arrange the text in the desired order
def arrange(array_1, array_2):
    array = []
    min_len = min(len(array_1), len(array_2)) # Determine the minimum length of the two arrays
    for i in range(min_len): # Loop through the minimum length
        array.append(array_1[i]) # Add the first element to the new array
        array.append(array_2[i]) # Add the second element to the new array
    return array # Return the new array

# Create a PowerPoint presentation
def createPptx(array, int_font_size, output_path):
    # font_size
    font_size = Pt(int_font_size)
    # Create a presentation
    prs = Presentation()
    # Set the slide layout and size
    slide_layout = prs.slide_layouts[6]
    prs.slide_width = Inches(13.33)
    prs.slide_heigth = Inches(4.8)
    # Calculate the center coordinates of the slide
    center_x = prs.slide_width / 2  
    # Calculate the size of the shape (e.g., a circle with a diameter of 2 inches)
    shape_width = Pt(690)

    for k in range(len(array)//2): # Use integer division to ensure the loop stops at the last pair
        # c_slide
        slide = prs.slides.add_slide(slide_layout)

        # Set background color black
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(0, 0, 0)

        # Set the standard of placeholders
        org_left_pt = center_x - shape_width / 2
        org_top_pt = Pt(150)
        org_width_pt = Pt(700)
        org_height_pt = Pt(157)

        # Organize the placeholders
        org_textbox = slide.shapes.add_textbox(left=org_left_pt, top=org_top_pt, width=org_width_pt, height=org_height_pt)
        org_text_frame = org_textbox.text_frame
        org_text_frame.word_wrap = True
        org_text_frame.text = array[2*k]  # Use 2k-1 as the index for the first element

        # Set font family, font size, and font color
        for paragraph in org_textbox.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = font_size
                run.font.name = 'Times New Roman'
                run.font.color.rgb = RGBColor(255, 255, 255)

        # Set the standard of trg_text
        trg_left_pt = center_x - shape_width / 2
        trg_top_pt = Inches(5.0)
        trg_width_pt = shape_width
        trg_height_pt = Pt(157)

        # Organize the placeholders
        trg_textbox = slide.shapes.add_textbox(left=trg_left_pt, top=trg_top_pt, width=trg_width_pt, height=trg_height_pt)
        trg_text_frame = trg_textbox.text_frame
        trg_text_frame.word_wrap = True
        trg_text_frame.text = array[2*k + 1]  # Use 2k as the index for the second element

        # Set font family, font size, and font color
        for paragraph in trg_textbox.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = font_size
                run.font.name = 'MS Mincho'
                run.font.color.rgb = RGBColor(255, 255, 255)

    # Save the presentation
    prs.save("../本番スライド/2024字幕.pptx")

# Specify the save location using a file dialog
output_path = specifySaveLocation()
# arrange
esja_arrange = arrange(es_text, ja_text)

# create .pptx
createPptx(esja_arrange, 40, output_path)